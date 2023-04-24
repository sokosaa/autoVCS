import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches
import comtypes.client
import os
import shutil
import time
import numpy
from PIL import Image
import pyautogui

import sys
from PyQt5.QtWidgets import QApplication, QWidget, QGridLayout, QPushButton, QFileDialog, QLabel, QLineEdit, QTextEdit, QMainWindow, QSizePolicy, QMessageBox

def interpret_path(path):
    path = path.text()
    path = path.split('\\')
    doubleslash_path = ''
    n = 0
    while len(path) > n:
        doubleslash_path = doubleslash_path + path[n]
        doubleslash_path = doubleslash_path + '\\\\'
        n += 1
    return doubleslash_path


def text2folder(working_path, part_numbers):
    try:
        working_path = interpret_path(working_path)
        list_of_part_numbers = part_numbers.toPlainText().splitlines()
        for line in list_of_part_numbers:
            line = line.strip()
            part_number = line
            # part_number = line.split(' ')[0]
            # vendor_code = line.split(' ')[1]
            if os.path.exists(working_path + part_number) == True:
                pass
            elif os.path.exists(working_path + part_number) == False:
                os.mkdir(working_path + part_number)
    except:
        display_error()


def rename_images(working_path, part_numbers):
    parent_path = interpret_path(working_path)
    list_of_part_numbers = part_numbers.toPlainText().splitlines()
    for line in list_of_part_numbers:
        line = line.strip()
        path3 = parent_path + line + '\\'
        images = os.listdir(path3)
        n = 0
        num_files = len(images)
        for i in range(num_files):
            while num_files > n:
                if path3+images[n] != line+'_VCS_img_'+str(n+1)+'.jpg' or path3+images[n] != line+'_VCS_img_'+str(n+1)+'.png':
                    os.rename(path3+images[n], path3+line+'_VCS_img_'+str(n+1)+'.jpg')
                    n += 1
                else:
                    n += 1


def create_ppt_function(working_path, part_numbers, template_pptx):
    # Check that the specified file exists
    if not os.path.isfile(template_pptx):
        raise FileNotFoundError(f"Template file '{template_pptx}' not found.")
    parent_path = interpret_path(working_path)
    list_of_part_numbers = part_numbers.toPlainText().splitlines()
    for line in list_of_part_numbers:
        line = line.strip()
        path2 = parent_path + line + '\\'
        images = os.listdir(path2)
        num_files = len(images)
        # if template_pptx == '':
        prs = Presentation(template_pptx)
        # else:
        #     prs = Presentation(interpret_path(template_pptx))
        slide1 = prs.slides[0]
        left = Inches(0)
        n = 0
        h = 3.75
        h_list = []
        h_array = numpy.array(h_list)
        while num_files > n:
            img = path2+images[n]
            j = Image.open(img)
            wide = j.width
            high = j.height
            ratio = (wide/high)
            top = numpy.sum(h_array)
            s = 0
            slide1.shapes.add_picture(img, left, top * 914400.1182, height=Inches(h))
            n += 1
        prs.save(parent_path+line+'_VCS.pptx')


def PPTtoPDF(working_path, part_numbers):
    parent_path = interpret_path(working_path)
    list_of_part_numbers = part_numbers.toPlainText().splitlines()
    powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
    powerpoint.Visible = 1
    for line in list_of_part_numbers:
        line = line.strip()
        part_number = line
        input_file = parent_path+part_number+'_VCS.pptx'
        output_file = parent_path+part_number+'_VCS.pdf'
        deck = powerpoint.Presentations.Open(input_file)
        deck.SaveAs(output_file, 32) # formatType = 32 for ppt to pdf. see: https://docs.microsoft.com/en-us/office/vba/api/powerpoint.ppsaveasfiletype
        deck.Close()


def move2focus(working_path, part_numbers, vendor_code, destinations_edit):
    parent_path = interpret_path(working_path)
    list_of_part_numbers = part_numbers.toPlainText().splitlines()
    for line in list_of_part_numbers:
        line = line.strip()
        source = parent_path + '\\' + line + '_VCS.pdf'
        destination = '\\\\focus\\InspPlanDocs\\' + line + '\\Incoming\\' + vendor_code
        existing_VCS = destination + '\\' + line + '_VCS.pdf'
        if os.path.exists(existing_VCS) == True:
            os.rename(existing_VCS, destination +
                        '\\' + line + '_VCS_old.pdf')
            destinations_edit.insertPlainText(destination+'\n')
        elif os.path.exists(destination) == True:
            pass
        else:
            destinations_edit.insertPlainText(destination+'\n')
            os.makedirs(destination)
        shutil.move(source, destination)


def RetractFromFocus():  # make this fucntion take files back from focus and into working path if user accidentally confirmed move to focus
    print('unfinished function RetractFromFocus')


def openDDS():  # IN#? of part no finds in # in dds then opens link with IN#? http://dds.rb.net/Forms/PPAP/L2L3PPAP.aspx?InspectionID=IN-259287
    print('unfinished function openDDS')


def enter_info():
    global vendor_code
    x = input('Paste working path: ')
    x = x.split('\\')
    path = ''
    n = 0
    while len(x) > n:
        path = path + x[n]
        path = path + '\\\\'
        n += 1

    file = input('Parts file name: ')
    if file[-4:] != '.txt':
        file = file + '.txt'
    file = path+file

    text2folder(path, file)
    print('Place images into their appropriate folders! Then answer the next question.')
    time.sleep(2)
    yesorno = input(
        'Are all parts from the same vendor? Enter "yes" or "no": ')
    if yesorno == 'yes' or yesorno == 'y':
        vendor_code = input('Vendor code: ')
    elif yesorno == 'no' or yesorno == 'n':
        print('Unable to move to focus, but you can find the VCS PDFs in your working path')
    else:
        print('invalid entry. program will close. please restart.')
        time.sleep(4)
        quit()

    rename_images(path, file)
    # create_ppt(path, file)
    input('Press enter if slides have been edited to will. Edited slides?')
    with open(file) as parts_text_file:
        for line in parts_text_file:
            line = line.strip()
            input_file = path + line + '_VCS.pptx'
            output_file = path + line + '_VCS.pdf'
            PPTtoPDF(input_file, output_file)
    with open(file) as parts_text_file:
        for line in parts_text_file:
            line = line.strip()
            input_file = path + line + '_VCS.pptx'
            os.remove(input_file)
    if yesorno == 'yes':
        while True:
            print('Please close PowerPoint. \n Please check the generated VCSs. \n If not using chrome as your pdf viewer, close the PDFs after you viewed them. \n Now, upload these VCSs to DDS. \n Now, would you like to move these VCSs to their appropriate focus folder? \n Enter "yes" or "no" below:')
            gomove = input()
            if gomove == 'yes':
                move2focus(path, file, vendor_code)
                break
            if gomove == 'no':
                print('Okay, you can delete the VCS PDFs if you do not like them or you can do whatever you want with them. \n Just keep in mind these files we be overwritten if you run this program again with the same part numbers.')
            print('invalid input, try again')


def display_error():

    # Create the message box
    msg_box = QMessageBox()

    # Set the message type to Error
    msg_box.setIcon(QMessageBox.Critical)

    # Set the text to be displayed
    msg_box.setText("An error has occurred, probably an invalid path or part number that has a character that can't be converted into a folder name. If its not that then ¯\_(ツ)_/¯")

    # Get the error message
    error_msg = str(sys.exc_info()[1])

    # Add the error message to the dialog box
    msg_box.setInformativeText(error_msg)

    # Display the message box
    msg_box.exec_()


if __name__ == '__main__':   
    pass

